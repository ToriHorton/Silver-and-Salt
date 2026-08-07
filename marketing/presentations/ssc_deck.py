#!/usr/bin/env python3
"""
Silver & Salt Capital — shared deck toolkit.

Brand tokens, colorway themes, slide helpers, and the 12-slide master template,
all built strictly from `_reference/brand-standards.md` and `BRAND.md`. Imported
by the individual build scripts:

  build_template.py         -> rust master template
  build_template_plum.py    -> plum master template
  build_pitch_deck.py       -> pre-filled pitch deck
  build_investor_update.py  -> pre-filled investor update

Google Slides imports the resulting .pptx natively (Drive: upload, then
Open with > Google Slides). Every element stays fully editable.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------
# Brand tokens (verbatim from the brand standards)
# --------------------------------------------------------------------------
CREAM      = RGBColor(0xFB, 0xF8, 0xF2)   # primary background
SAND       = RGBColor(0xF4, 0xEF, 0xE6)   # secondary background, cards
MOSS       = RGBColor(0x2F, 0x3E, 0x34)   # primary text / dark surfaces
MOSS_LIGHT = RGBColor(0x4A, 0x5E, 0x50)   # body text
SAGE       = RGBColor(0x7E, 0x8E, 0x84)   # muted text, labels, tagline
RUST       = RGBColor(0xD1, 0x6B, 0x4F)   # rust accent
RUST_HOVER = RGBColor(0xE0, 0x7A, 0x5E)   # rust highlight
PLUM       = RGBColor(0x8B, 0x5E, 0x83)   # plum accent (How It Works palette)
WARM       = RGBColor(0xC4, 0xA4, 0x7E)   # gold accent
CREAM_TEXT = RGBColor(0xF0, 0xED, 0xE7)   # cream wordmark on dark
HAIRLINE   = RGBColor(0xE3, 0xDD, 0xD1)   # faint rule on cream

# Pre-blended watermark colors (baked "whisper" so it reads the same in any
# renderer that lacks true text alpha).
WM_ON_CREAM = RGBColor(0xF3, 0xF1, 0xEA)  # moss @ ~4% over cream
WM_ON_SAND  = RGBColor(0xEA, 0xE6, 0xDD)  # moss @ ~5% over sand
WM_ON_MOSS  = RGBColor(0x3C, 0x4A, 0x40)  # white @ ~6% over moss

# Fonts (all three exist in Google Slides' "More fonts" picker).
SERIF   = "Cormorant Garamond"   # display, headings, wordmark, quotes
NUMERAL = "Cormorant Infant"     # digits only (flagged "1")
SANS    = "Hanken Grotesk"       # body, labels, UI (Satoshi substitute)

# 16:9 canvas
SW = Inches(13.333)
SH = Inches(7.5)
MARGIN = Inches(0.92)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "..", "logos", "brand-assets")


def asset(name):
    return os.path.join(ASSETS, name)


ICON = asset("ssc-icon.png")
ICON_PLUM = asset("ssc-icon-plum.png")


# --------------------------------------------------------------------------
# Colorway theme
# --------------------------------------------------------------------------
class Theme:
    def __init__(self, name, accent, accent_on_dark, divider_bg, divider_wm,
                 divider_sub, logo_light, logo_dark, icon, extra_swatches):
        self.name = name
        self.accent = accent                 # accent on light surfaces
        self.accent_on_dark = accent_on_dark # accent text on moss
        self.divider_bg = divider_bg         # full-bleed divider surface
        self.divider_wm = divider_wm         # watermark on divider
        self.divider_sub = divider_sub       # subtitle text on divider
        self.logo_light = logo_light         # for dark backgrounds
        self.logo_dark = logo_dark           # for light backgrounds
        self.icon = icon                     # footer badge
        self.extra_swatches = extra_swatches # appendix palette accents


RUST_THEME = Theme(
    name="Rust",
    accent=RUST, accent_on_dark=RUST,
    divider_bg=RUST, divider_wm=RGBColor(0xD9, 0x7B, 0x61),
    divider_sub=RGBColor(0xF5, 0xE2, 0xDA),
    logo_light=asset("ssc-primary-logo-light.png"),
    logo_dark=asset("ssc-primary-logo.png"),
    icon=ICON,
    extra_swatches=[("Rust", RUST, "#D16B4F", CREAM_TEXT),
                    ("Rust Hover", RUST_HOVER, "#E07A5E", MOSS)],
)

PLUM_THEME = Theme(
    name="Plum",
    accent=PLUM, accent_on_dark=RGBColor(0xB5, 0x89, 0xAC),
    divider_bg=PLUM, divider_wm=RGBColor(0x99, 0x71, 0x92),
    divider_sub=RGBColor(0xED, 0xE6, 0xEC),
    logo_light=asset("ssc-primary-plum-light.png"),
    logo_dark=asset("ssc-primary-plum.png"),
    icon=ICON_PLUM,
    extra_swatches=[("Plum", PLUM, "#8B5E83", CREAM_TEXT),
                    ("Plum Light", RGBColor(0xB5, 0x89, 0xAC), "#B589AC",
                     MOSS)],
)


# --------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    set_fill(r, color)
    r.shadow.inherit = False
    return r


def rule(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(r, color)
    r.shadow.inherit = False
    return r


def _tracking(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        run = p.add_run()
        run.text = ln["text"]
        f = run.font
        f.name = ln.get("font", SANS)
        f.size = Pt(ln["size"])
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.color.rgb = ln.get("color", MOSS)
        if ln.get("tracking") is not None:
            _tracking(run, ln["tracking"])
    return tb


def watermark(slide, color, *, x=Inches(7.2), size=560):
    tb = slide.shapes.add_textbox(x, Inches(-1.2), Inches(7.0), Inches(10.0))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "&"
    run.font.name = SERIF
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return tb


def picture(slide, path, x, y, *, height=None, width=None):
    if height is not None:
        return slide.shapes.add_picture(path, x, y, height=height)
    return slide.shapes.add_picture(path, x, y, width=width)


def eyebrow(slide, text, x, y, color, w=Inches(6)):
    return textbox(slide, x, y, w, Inches(0.4),
                   [{"text": text.upper(), "font": SANS, "size": 12,
                     "color": color, "bold": True, "tracking": 2.6}])


def footer(slide, page, total, theme, *, label="Silver & Salt Capital"):
    picture(slide, theme.icon, MARGIN, Inches(6.92), height=Inches(0.28))
    textbox(slide, MARGIN + Inches(0.42), Inches(6.93), Inches(5.0),
            Inches(0.4),
            [{"text": label, "font": SERIF, "size": 14, "color": SAGE,
              "tracking": 0.2}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, SW - MARGIN - Inches(2.0), Inches(6.93), Inches(2.0),
            Inches(0.4),
            [{"text": f"{page:02d} / {total:02d}", "font": NUMERAL,
              "size": 13, "color": SAGE, "tracking": 0.6,
              "align": PP_ALIGN.RIGHT}],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=SAND, radius=0.05):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(c, fill)
    c.shadow.inherit = False
    c.adjustments[0] = radius
    return c


# --------------------------------------------------------------------------
# Composite slide builders (reused by templates and example decks)
# --------------------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def add_slide(prs, fill):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, fill)
    return s


def cover(prs, theme, *, title, subtitle, tagline=None, eyebrow_text=None,
          date_text=None):
    tagline = tagline or "Connecting capital to Utah founders who use it best."
    s = add_slide(prs, MOSS)
    watermark(s, WM_ON_MOSS, x=Inches(7.6), size=620)
    picture(s, theme.logo_light, MARGIN, Inches(0.9), height=Inches(0.62))
    if eyebrow_text:
        eyebrow(s, eyebrow_text, MARGIN, Inches(2.15), theme.accent_on_dark)
    rule(s, MARGIN, Inches(2.55), Inches(0.9), Pt(2.2), theme.accent_on_dark)
    textbox(s, MARGIN, Inches(2.75), Inches(11.0), Inches(2.6),
            [{"text": title, "font": SERIF, "size": 68, "color": CREAM_TEXT,
              "tracking": -0.4, "line_spacing": 1.02, "space_after": 6},
             {"text": subtitle, "font": SERIF, "size": 30, "color": SAGE,
              "italic": True, "line_spacing": 1.1}])
    textbox(s, MARGIN, Inches(5.7), Inches(9.0), Inches(0.6),
            [{"text": tagline, "font": SERIF, "size": 19, "color": SAGE,
              "italic": True}])
    if date_text:
        textbox(s, SW - MARGIN - Inches(3.4), Inches(5.72), Inches(3.4),
                Inches(0.5),
                [{"text": date_text.upper(), "font": SANS, "size": 12,
                  "color": SAGE, "bold": True, "tracking": 2.4,
                  "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)
    return s


def section_divider(prs, theme, *, number, title, subtitle, on_moss=False):
    """Divider on the accent surface (on_moss=False) or on moss (True)."""
    if on_moss:
        s = add_slide(prs, MOSS)
        watermark(s, WM_ON_MOSS, x=Inches(7.8), size=640)
        num_color = theme.accent_on_dark
        title_color = CREAM_TEXT
        sub_color = SAGE
        rule_color = theme.accent_on_dark
    else:
        s = add_slide(prs, theme.divider_bg)
        watermark(s, theme.divider_wm, x=Inches(7.8), size=640)
        num_color = CREAM_TEXT
        title_color = CREAM_TEXT
        sub_color = theme.divider_sub
        rule_color = CREAM_TEXT
    textbox(s, MARGIN, Inches(2.5), Inches(3.0), Inches(2.0),
            [{"text": number, "font": NUMERAL, "size": 150, "color": num_color,
              "tracking": -2, "line_spacing": 1.0}])
    rule(s, MARGIN, Inches(4.55), Inches(0.9), Pt(2.2), rule_color)
    lines = [{"text": title, "font": SERIF, "size": 56, "color": title_color,
              "tracking": -0.4, "space_after": 4}]
    if subtitle:
        lines.append({"text": subtitle, "font": SERIF, "size": 24,
                      "color": sub_color, "italic": True})
    textbox(s, MARGIN, Inches(4.75), Inches(11), Inches(1.4), lines)
    return s


def content(prs, theme, *, eyebrow_text, heading, body, page, total,
            callout=None, callout_source=None):
    """Heading + body column, with an optional right-hand callout panel."""
    s = add_slide(prs, CREAM)
    eyebrow(s, eyebrow_text, MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(11.4), Inches(1.5),
            [{"text": heading, "font": SERIF, "size": 44, "color": MOSS,
              "tracking": -0.4, "line_spacing": 1.02}])
    rule(s, MARGIN, Inches(2.85), Inches(0.9), Pt(2.2), theme.accent)
    body_w = Inches(7.0) if callout else Inches(11.0)
    textbox(s, MARGIN, Inches(3.25), body_w, Inches(3.2),
            [{"text": b, "font": SANS, "size": 18, "color": MOSS_LIGHT,
              "line_spacing": 1.5, "space_after": 14} for b in body])
    if callout:
        card(s, Inches(8.4), Inches(3.25), Inches(4.0), Inches(3.0),
             radius=0.04)
        lines = [{"text": callout, "font": SERIF, "size": 24, "color": MOSS,
                  "italic": True, "line_spacing": 1.1, "space_after": 10}]
        if callout_source:
            lines.append({"text": callout_source, "font": SANS, "size": 13,
                          "color": SAGE, "tracking": 0.4})
        textbox(s, Inches(8.8), Inches(3.55), Inches(3.3), Inches(2.5), lines,
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total, theme)
    return s


def two_column(prs, theme, *, eyebrow_text, heading, columns, page, total):
    """columns: list of exactly 2 dicts {head, body}."""
    s = add_slide(prs, CREAM)
    eyebrow(s, eyebrow_text, MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(11), Inches(1.0),
            [{"text": heading, "font": SERIF, "size": 44, "color": MOSS,
              "tracking": -0.4}])
    accents = [theme.accent, MOSS]
    xs = [Inches(0.92), Inches(6.95)]
    for i, col in enumerate(columns[:2]):
        x = xs[i]
        rule(s, x, Inches(3.0), Inches(0.7), Pt(2.2), accents[i])
        textbox(s, x, Inches(3.2), Inches(5.2), Inches(0.7),
                [{"text": col["head"], "font": SERIF, "size": 30,
                  "color": MOSS}])
        textbox(s, x, Inches(3.95), Inches(5.2), Inches(2.6),
                [{"text": col["body"], "font": SANS, "size": 17,
                  "color": MOSS_LIGHT, "line_spacing": 1.5}])
    footer(s, page, total, theme)
    return s


def stats(prs, theme, *, eyebrow_text, heading, items, page, total,
          surface=SAND):
    """items: list of up to 3 dicts {num, label, accent(optional)}."""
    s = add_slide(prs, surface)
    wm = WM_ON_SAND if surface == SAND else WM_ON_CREAM
    watermark(s, wm, x=Inches(10.4))
    eyebrow(s, eyebrow_text, MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(11), Inches(1.0),
            [{"text": heading, "font": SERIF, "size": 44, "color": MOSS,
              "tracking": -0.4}])
    col_w = Inches(3.6)
    gap = Inches(0.4)
    for i, it in enumerate(items[:3]):
        accent = it.get("accent", theme.accent if i != 1 else MOSS)
        x = MARGIN + (col_w + gap) * i
        textbox(s, x, Inches(3.0), col_w, Inches(1.7),
                [{"text": it["num"], "font": NUMERAL, "size": 108,
                  "color": accent, "tracking": -2, "line_spacing": 0.95}])
        rule(s, x + Inches(0.08), Inches(4.75), Inches(0.6), Pt(2), accent)
        textbox(s, x, Inches(4.95), col_w, Inches(1.2),
                [{"text": it["label"], "font": SANS, "size": 16,
                  "color": MOSS_LIGHT, "line_spacing": 1.4}])
    footer(s, page, total, theme)
    return s


def pull_quote(prs, theme, *, quote, attribution):
    s = add_slide(prs, MOSS)
    watermark(s, WM_ON_MOSS, x=Inches(7.8), size=640)
    textbox(s, MARGIN, Inches(1.6), Inches(2.0), Inches(2.0),
            [{"text": "“", "font": SERIF, "size": 200,
              "color": theme.accent_on_dark, "line_spacing": 1.0}])
    textbox(s, MARGIN, Inches(2.95), Inches(10.4), Inches(2.8),
            [{"text": quote, "font": SERIF, "size": 42, "color": CREAM_TEXT,
              "italic": True, "line_spacing": 1.12}])
    rule(s, MARGIN, Inches(5.95), Inches(0.5), Pt(2), theme.accent_on_dark)
    textbox(s, MARGIN + Inches(0.7), Inches(5.78), Inches(9), Inches(0.5),
            [{"text": attribution, "font": SANS, "size": 15, "color": SAGE,
              "tracking": 1.2}])
    return s


def cards_grid(prs, theme, *, eyebrow_text, heading, cards, page, total):
    """cards: list of up to 3 dicts {num, head, body}."""
    s = add_slide(prs, CREAM)
    eyebrow(s, eyebrow_text, MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(11), Inches(1.0),
            [{"text": heading, "font": SERIF, "size": 44, "color": MOSS,
              "tracking": -0.4}])
    cw = Inches(3.6)
    cgap = Inches(0.4)
    for i, c in enumerate(cards[:3]):
        x = MARGIN + (cw + cgap) * i
        card(s, x, Inches(2.95), cw, Inches(3.4))
        textbox(s, x + Inches(0.35), Inches(3.25), Inches(2.0), Inches(0.8),
                [{"text": c["num"], "font": NUMERAL, "size": 34,
                  "color": theme.accent}])
        textbox(s, x + Inches(0.35), Inches(4.05), cw - Inches(0.7),
                Inches(0.7),
                [{"text": c["head"], "font": SERIF, "size": 26,
                  "color": MOSS}])
        textbox(s, x + Inches(0.35), Inches(4.75), cw - Inches(0.7),
                Inches(1.5),
                [{"text": c["body"], "font": SANS, "size": 15,
                  "color": MOSS_LIGHT, "line_spacing": 1.45}])
    footer(s, page, total, theme)
    return s


def agenda(prs, theme, *, heading, items, page, total):
    """items: list of (number, title, desc)."""
    s = add_slide(prs, CREAM)
    watermark(s, WM_ON_CREAM)
    eyebrow(s, "Contents", MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(8), Inches(1.0),
            [{"text": heading, "font": SERIF, "size": 46, "color": MOSS,
              "tracking": -0.4}])
    top = Inches(2.55)
    row_h = Inches(0.86)
    for i, (num, title, desc) in enumerate(items):
        y = top + row_h * i
        textbox(s, MARGIN, y, Inches(0.95), row_h,
                [{"text": num, "font": NUMERAL, "size": 30,
                  "color": theme.accent}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, MARGIN + Inches(1.15), y, Inches(4.2), row_h,
                [{"text": title, "font": SERIF, "size": 27, "color": MOSS}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, MARGIN + Inches(5.4), y, Inches(5.4), row_h,
                [{"text": desc, "font": SANS, "size": 15,
                  "color": MOSS_LIGHT}], anchor=MSO_ANCHOR.MIDDLE)
        rule(s, MARGIN, y + row_h - Inches(0.02), Inches(10.5), Pt(0.75),
             HAIRLINE)
    footer(s, page, total, theme)
    return s


def closing(prs, theme, *, line1, line2, url="silverandsaltcapital.com",
            email="tori@silverandsaltcapital.com", cta="Join Us"):
    s = add_slide(prs, MOSS)
    watermark(s, WM_ON_MOSS, x=Inches(7.6), size=620)
    textbox(s, MARGIN, Inches(1.7), Inches(11.5), Inches(2.6),
            [{"text": line1, "font": SERIF, "size": 46, "color": CREAM_TEXT,
              "bold": True, "tracking": -0.3, "line_spacing": 1.08,
              "space_after": 8},
             {"text": line2, "font": SERIF, "size": 46,
              "color": theme.accent_on_dark, "bold": True, "tracking": -0.3,
              "line_spacing": 1.08}])
    # On moss the wordmark must read in cream: use the on-dark light logo.
    picture(s, theme.logo_light, MARGIN, Inches(5.05), height=Inches(0.52))
    textbox(s, SW - MARGIN - Inches(4.6), Inches(5.0), Inches(4.6),
            Inches(1.6),
            [{"text": cta, "font": SANS, "size": 16,
              "color": theme.accent_on_dark, "bold": True, "tracking": 1.4,
              "align": PP_ALIGN.RIGHT, "space_after": 6},
             {"text": url, "font": SANS, "size": 15, "color": CREAM_TEXT,
              "align": PP_ALIGN.RIGHT, "space_after": 2},
             {"text": email, "font": SANS, "size": 15, "color": SAGE,
              "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)
    return s


def appendix_palette(prs, theme, *, page, total):
    s = add_slide(prs, CREAM)
    eyebrow(s, "Appendix", MARGIN, Inches(0.95), theme.accent)
    textbox(s, MARGIN, Inches(1.35), Inches(11), Inches(1.0),
            [{"text": "Palette & type", "font": SERIF, "size": 44,
              "color": MOSS, "tracking": -0.4}])
    swatches = [("Cream", CREAM, "#FBF8F2", MOSS),
                ("Sand", SAND, "#F4EFE6", MOSS),
                ("Moss", MOSS, "#2F3E34", CREAM_TEXT),
                ("Moss Light", MOSS_LIGHT, "#4A5E50", CREAM_TEXT),
                ("Sage", SAGE, "#7E8E84", CREAM_TEXT)]
    swatches += theme.extra_swatches
    swatches += [("Warm", WARM, "#C4A47E", MOSS)]
    sw_w = Inches(1.35)
    sw_gap = Inches(0.1)
    for i, (nm, col, hexv, txt) in enumerate(swatches):
        x = MARGIN + (sw_w + sw_gap) * i
        chip = card(s, x, Inches(2.7), sw_w, Inches(1.5), fill=col,
                    radius=0.08)
        chip.line.color.rgb = HAIRLINE
        chip.line.width = Pt(0.75)
        textbox(s, x + Inches(0.12), Inches(3.45), sw_w - Inches(0.2),
                Inches(0.7),
                [{"text": nm, "font": SANS, "size": 11, "color": txt,
                  "bold": True, "space_after": 1},
                 {"text": hexv, "font": NUMERAL, "size": 11, "color": txt}])
    textbox(s, MARGIN, Inches(4.7), Inches(11), Inches(0.8),
            [{"text": "Cormorant Garamond  —  display, headings, quotes",
              "font": SERIF, "size": 30, "color": MOSS}])
    textbox(s, MARGIN, Inches(5.45), Inches(11), Inches(0.6),
            [{"text": "Cormorant Infant  /  1234567890  —  numerals only "
              "(flagged 1)", "font": NUMERAL, "size": 24,
              "color": theme.accent}])
    textbox(s, MARGIN, Inches(6.05), Inches(11), Inches(0.6),
            [{"text": "Hanken Grotesk  —  body, labels, and UI "
              "(Satoshi substitute in Google Slides)", "font": SANS,
              "size": 17, "color": MOSS_LIGHT}])
    footer(s, page, total, theme)
    return s


def save(prs, filename):
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), filename)
    prs.save(out)
    print("Wrote", filename, "-", len(prs.slides._sldIdLst), "slides")
    return out


# --------------------------------------------------------------------------
# The 12-slide master template (themed)
# --------------------------------------------------------------------------
def build_template(theme, filename):
    T = 12
    prs = new_deck()
    cover(prs, theme, title="Presentation Title",
          subtitle="A calm, declarative subtitle goes on this second line.",
          date_text="Month 2026")
    agenda(prs, theme, heading="What we will cover", items=[
        ("01", "The opportunity", "Why Utah, why now."),
        ("02", "Our thesis", "Where capital meets conviction."),
        ("03", "How it works", "The model, in plain terms."),
        ("04", "The numbers", "What the data shows."),
        ("05", "Join us", "Where you fit, and what comes next."),
    ], page=2, total=T)
    section_divider(prs, theme, number="01", title="Section Title",
                    subtitle="One line that frames what this section is "
                    "about.")
    content(prs, theme, eyebrow_text="The Opportunity",
            heading="A heading that states the point directly",
            body=[
                "Lead with the idea in one calm sentence. Short, declarative, "
                "confident.",
                "Use the body style for supporting detail. Generous spacing, "
                "easy to read from the back of the room.",
                "Each point earns its place. When a line wants a dash, it "
                "usually wants to be two sentences instead.",
            ],
            callout="A supporting stat or callout sits here.",
            callout_source="Source name", page=4, total=T)
    two_column(prs, theme, eyebrow_text="How It Works",
               heading="Two ideas, side by side", columns=[
                   {"head": "Column one",
                    "body": "Open with the positive. State what this is, "
                    "plainly, then give two or three lines of supporting "
                    "detail underneath."},
                   {"head": "Column two",
                    "body": "Mirror the structure on the right. Balanced "
                    "columns read as calm and considered, which is exactly "
                    "the brand voice."},
               ], page=5, total=T)
    stats(prs, theme, eyebrow_text="The Numbers",
          heading="What the data shows", items=[
              {"num": "#3", "label": "Utah's rank to start a business"},
              {"num": "#50", "label": "Utah's rank for women's equality"},
              {"num": "169", "label": "An example headline figure"},
          ], page=6, total=T)
    pull_quote(prs, theme,
               quote="A pull quote sits here, set large and italic in "
               "Cormorant Garamond. It carries weight because it is given "
               "room to breathe.",
               attribution="Attribution, Title")
    # Image + caption (kept inline; not a common reusable signature)
    s = add_slide(prs, CREAM)
    img = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.4), SH)
    set_fill(img, SAND)
    img.shadow.inherit = False
    textbox(s, Inches(0.5), Inches(3.3), Inches(5.4), Inches(1.0),
            [{"text": "Replace with image", "font": SANS, "size": 15,
              "color": SAGE, "tracking": 1.0, "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)
    eyebrow(s, "Story", Inches(7.05), Inches(1.4), theme.accent)
    textbox(s, Inches(7.05), Inches(1.8), Inches(5.4), Inches(1.6),
            [{"text": "A picture, then the point", "font": SERIF, "size": 40,
              "color": MOSS, "tracking": -0.4, "line_spacing": 1.02}])
    rule(s, Inches(7.05), Inches(3.5), Inches(0.7), Pt(2.2), theme.accent)
    textbox(s, Inches(7.05), Inches(3.85), Inches(5.4), Inches(2.6),
            [{"text": "Pair a full-bleed image on the left with a tight "
              "column of copy on the right. Let the image do the emotional "
              "work and keep the words quiet underneath it.", "font": SANS,
              "size": 17, "color": MOSS_LIGHT, "line_spacing": 1.5}])
    footer(s, 8, T, theme)
    cards_grid(prs, theme, eyebrow_text="The Collective",
               heading="Three of something", cards=[
                   {"num": "01", "head": "First card",
                    "body": "A short description that explains this one item "
                    "in two calm lines."},
                   {"num": "02", "head": "Second card",
                    "body": "Keep each card parallel in length so the row "
                    "reads as one considered unit."},
                   {"num": "03", "head": "Third card",
                    "body": "End on the strongest of the three. Define each "
                    "by what it is."},
               ], page=9, total=T)
    section_divider(prs, theme, number="02", title="Second Section",
                    subtitle="An alternate divider on the moss surface.",
                    on_moss=True)
    closing(prs, theme,
            line1="Utah's next generation of innovation will be shaped by "
            "founders & determined by funders.",
            line2="It's time more of them are women.")
    appendix_palette(prs, theme, page=12, total=T)
    return save(prs, filename)
